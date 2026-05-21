#!/usr/bin/env python3
"""
答辩 PPT：AI 驱动的靶场自动生成系统
简洁深色科技风 · 13 张 · 16:9 宽屏
"""
import sys
sys.path.insert(0, '/Users/gw/Library/Python/3.9/lib/python/site-packages')

from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.enum.shapes import MSO_AUTO_SHAPE_TYPE as MSO
from pptx.oxml.ns import qn
from lxml import etree

# ─── 色板 ────────────────────────────────────────────────────────────────────
BG      = RGBColor(0x0D, 0x12, 0x27)   # 深蓝黑
PANEL   = RGBColor(0x14, 0x1D, 0x3B)   # 卡片底色
PANEL2  = RGBColor(0x1B, 0x26, 0x4A)   # 高亮卡片
LINE    = RGBColor(0x1E, 0x2D, 0x55)   # 细线/分隔
CYAN    = RGBColor(0x00, 0xC8, 0xE0)   # 主色-青
BLUE    = RGBColor(0x38, 0x8B, 0xFF)   # 辅色-蓝
PURPLE  = RGBColor(0x8B, 0x5C, 0xF6)   # 辅色-紫
GREEN   = RGBColor(0x3D, 0xD6, 0x8C)   # 成功-绿
ORANGE  = RGBColor(0xFF, 0xA5, 0x3B)   # 警告-橙
RED     = RGBColor(0xFF, 0x5B, 0x5B)   # 错误-红
WHITE   = RGBColor(0xFF, 0xFF, 0xFF)
SILVER  = RGBColor(0xC8, 0xD4, 0xE8)   # 正文
MUTED   = RGBColor(0x6B, 0x7A, 0x99)   # 次要文字
ACCENT  = RGBColor(0x00, 0x80, 0xA8)   # 青-暗（边框）

W = Inches(13.33)
H = Inches(7.5)


# ─── 工具 ────────────────────────────────────────────────────────────────────
def darken(c, div=4):
    s = str(c)
    return RGBColor(int(s[0:2],16)//div, int(s[2:4],16)//div, int(s[4:6],16)//div)

def lighten(c, add=30):
    s = str(c)
    def clamp(v): return min(255, v + add)
    return RGBColor(clamp(int(s[0:2],16)), clamp(int(s[2:4],16)), clamp(int(s[4:6],16)))


prs = Presentation()
prs.slide_width  = W
prs.slide_height = H
BLANK = prs.slide_layouts[6]


def new_slide():
    s = prs.slides.add_slide(BLANK)
    s.background.fill.solid()
    s.background.fill.fore_color.rgb = BG
    return s


def R(s, x, y, w, h, fill=PANEL, border=LINE, bw=1.0,
       shape=MSO.RECTANGLE, radius=None):
    """矩形/圆角矩形"""
    sh = s.shapes.add_shape(shape, x, y, w, h)
    sh.fill.solid()
    sh.fill.fore_color.rgb = fill
    if bw > 0:
        sh.line.color.rgb = border
        sh.line.width = Pt(bw)
    else:
        sh.line.fill.background()
    return sh


def T(s, text, x, y, w, h, size=13, bold=False, color=SILVER,
      align=PP_ALIGN.LEFT, italic=False, wrap=True):
    """文本框"""
    tb = s.shapes.add_textbox(x, y, w, h)
    tf = tb.text_frame
    tf.word_wrap = wrap
    p = tf.paragraphs[0]
    p.alignment = align
    r = p.add_run()
    r.text = text
    r.font.size = Pt(size)
    r.font.bold = bold
    r.font.italic = italic
    r.font.color.rgb = color
    return tb


def hline(s, y, x=Inches(0.5), w=Inches(12.33), color=LINE):
    R(s, x, y, w, Inches(0.015), fill=color, border=color, bw=0)


def slide_header(s, title, sub=''):
    """顶部标题区（带左侧色块+分隔线）"""
    # 顶部细彩条
    R(s, 0, 0, W, Inches(0.06), fill=CYAN, border=CYAN, bw=0)
    # 左侧竖条
    R(s, Inches(0.35), Inches(0.1), Inches(0.06), Inches(0.75),
      fill=CYAN, border=CYAN, bw=0)
    # 标题
    T(s, title, Inches(0.52), Inches(0.1), Inches(12.3), Inches(0.58),
      size=26, bold=True, color=WHITE)
    # 副标题
    if sub:
        T(s, sub, Inches(0.52), Inches(0.68), Inches(12.3), Inches(0.32),
          size=11, color=MUTED)
    # 分隔线
    hline(s, Inches(1.05), Inches(0.35), Inches(12.6), ACCENT)


def card(s, x, y, w, h, title='', lines=None,
         tc=CYAN, bc=ACCENT, fc=PANEL):
    """卡片：可选标题 + 多行正文"""
    R(s, x, y, w, h, fill=fc, border=bc, bw=1.0,
      shape=MSO.ROUNDED_RECTANGLE)
    # 左侧色条
    R(s, x, y, Inches(0.05), h, fill=tc, border=tc, bw=0,
      shape=MSO.ROUNDED_RECTANGLE)
    cy = y + Inches(0.12)
    if title:
        T(s, title, x+Inches(0.13), cy, w-Inches(0.22), Inches(0.38),
          size=12, bold=True, color=tc)
        cy += Inches(0.38)
    if lines:
        for ln in lines:
            T(s, ln, x+Inches(0.13), cy, w-Inches(0.22), Inches(0.32),
              size=10.5, color=SILVER)
            cy += Inches(0.3)


def stat_box(s, x, y, w, h, num, unit, label, color=CYAN):
    """大数字统计框"""
    R(s, x, y, w, h, fill=PANEL2, border=color, bw=1.2,
      shape=MSO.ROUNDED_RECTANGLE)
    T(s, num,   x, y+Inches(0.18), w, Inches(0.65),
      size=36, bold=True, color=color, align=PP_ALIGN.CENTER)
    T(s, unit,  x, y+Inches(0.72), w, Inches(0.28),
      size=11, bold=True, color=color, align=PP_ALIGN.CENTER)
    T(s, label, x, y+Inches(0.98), w, Inches(0.28),
      size=10, color=MUTED, align=PP_ALIGN.CENTER)


def accent_list(s, items, x, y, w, gap=Inches(0.35), color=CYAN, size=11):
    """带前缀色块的列表"""
    for i, item in enumerate(items):
        cy = y + gap * i
        R(s, x, cy+Inches(0.06), Inches(0.06), Inches(0.2),
          fill=color, border=color, bw=0)
        T(s, item, x+Inches(0.12), cy, w-Inches(0.12), Inches(0.32),
          size=size, color=SILVER)


def flow_arrow(s, x, y, w=Inches(0.35), h=Inches(0.25), color=MUTED):
    """向右箭头"""
    R(s, x, y, w, h, fill=color, border=color, bw=0,
      shape=MSO.RIGHT_ARROW)


def table_row(s, x, y, w, h, cells, widths, colors,
              fill=PANEL, text_color=SILVER, size=10.5, bold=False):
    cx = x
    for i, (cell, cw) in enumerate(zip(cells, widths)):
        bg = colors[i] if colors else fill
        R(s, cx, y, cw, h, fill=bg, border=LINE, bw=0.5)
        T(s, cell, cx+Inches(0.08), y+Inches(0.04), cw-Inches(0.16), h-Inches(0.08),
          size=size, color=text_color, bold=bold, align=PP_ALIGN.CENTER)
        cx += cw


# ═══════════════════════════════════════════════════════════════════════════════
# Slide 1  封面
# ═══════════════════════════════════════════════════════════════════════════════
s = new_slide()

# 背景渐变感：左侧暗色矩形
R(s, 0, 0, Inches(5), H, fill=RGBColor(0x08,0x0C,0x1E), border=BG, bw=0)

# 右侧装饰色块
R(s, Inches(11.5), 0, Inches(1.83), H,
  fill=RGBColor(0x00,0x30,0x48), border=BG, bw=0)
R(s, Inches(12.5), 0, Inches(0.83), H,
  fill=RGBColor(0x00,0x18,0x28), border=BG, bw=0)

# 顶部彩条
R(s, 0, 0, W, Inches(0.08), fill=CYAN, border=CYAN, bw=0)
R(s, 0, Inches(7.42), W, Inches(0.08), fill=PURPLE, border=PURPLE, bw=0)

# 主标题区
R(s, Inches(0.8), Inches(1.8), Inches(10.5), Inches(3.7),
  fill=RGBColor(0x10,0x18,0x32), border=CYAN, bw=2,
  shape=MSO.ROUNDED_RECTANGLE)

T(s, 'AI 驱动的靶场自动生成系统',
  Inches(1.2), Inches(2.1), Inches(9.7), Inches(1.0),
  size=36, bold=True, color=CYAN, align=PP_ALIGN.CENTER)

T(s, 'AI-Driven Cybersecurity Range Auto-Generation System',
  Inches(1.2), Inches(3.05), Inches(9.7), Inches(0.42),
  size=13, italic=True, color=SILVER, align=PP_ALIGN.CENTER)

hline(s, Inches(3.55), Inches(2.5), Inches(7.7), ACCENT)

T(s, '计算机科学与技术专业  ·  本科毕业设计',
  Inches(1.2), Inches(3.65), Inches(9.7), Inches(0.38),
  size=12, color=MUTED, align=PP_ALIGN.CENTER)

T(s, '汇报人：XXX          指导教师：XXX          2026 年 5 月',
  Inches(1.2), Inches(4.1), Inches(9.7), Inches(0.35),
  size=11, color=MUTED, align=PP_ALIGN.CENTER)

# 技术标签行
tags   = ['Flask + Vue 3', 'Docker SDK', 'DeepSeek LLM', 'Multi-Agent']
tcolors = [CYAN, GREEN, ORANGE, PURPLE]
for i, (tag, tc) in enumerate(zip(tags, tcolors)):
    bx = Inches(1.4) + Inches(2.55) * i
    R(s, bx, Inches(5.75), Inches(2.25), Inches(0.38),
      fill=darken(tc, 6), border=tc, bw=1.0,
      shape=MSO.ROUNDED_RECTANGLE)
    T(s, tag, bx, Inches(5.75), Inches(2.25), Inches(0.38),
      size=10.5, bold=True, color=tc, align=PP_ALIGN.CENTER)


# ═══════════════════════════════════════════════════════════════════════════════
# Slide 2  研究背景
# ═══════════════════════════════════════════════════════════════════════════════
s = new_slide()
slide_header(s, '研究背景：为什么需要 AI 靶场',
             '奇安信2024年度报告：网络安全威胁复杂度持续上升，实战演训环境亟待提升')

# 三个痛点卡
problems = [
    ('部署周期长',  ['传统靶场从需求到上线', '通常需要 数天至数周', '严重制约演练节奏']),
    ('场景高度固化', ['配置一旦完成难以更改', '不同目标需重复建设', '缺乏动态适应能力']),
    ('依赖人工运维', ['需熟悉 Docker / 网络 / 漏洞', '人工成本高、易出错', '门槛高阻碍小团队']),
]
for i, (hdr, lines) in enumerate(problems):
    cx = Inches(0.4) + Inches(4.31) * i
    card(s, cx, Inches(1.18), Inches(4.0), Inches(3.5),
         title=hdr, lines=lines, tc=RED, bc=RGBColor(0x55,0x15,0x15),
         fc=RGBColor(0x1A,0x10,0x1A))

# 解决方案框
R(s, Inches(0.4), Inches(5.0), Inches(12.53), Inches(1.95),
  fill=RGBColor(0x05,0x20,0x35), border=CYAN, bw=1.5,
  shape=MSO.ROUNDED_RECTANGLE)
R(s, Inches(0.4), Inches(5.0), Inches(0.07), Inches(1.95),
  fill=CYAN, border=CYAN, bw=0, shape=MSO.ROUNDED_RECTANGLE)
T(s, '本文方案',
  Inches(0.58), Inches(5.1), Inches(2.5), Inches(0.4),
  size=13, bold=True, color=CYAN)
T(s, '自然语言描述  -->  LLM 解析意图  -->  JSON 配置  -->  Docker 自动编排部署',
  Inches(0.58), Inches(5.5), Inches(12.0), Inches(0.38),
  size=12, bold=True, color=WHITE)
T(s, '目标：分钟级完成靶场搭建 · 用户无需具备 Docker/网络配置知识 · 适用轻量级教学实验与攻防演练',
  Inches(0.58), Inches(5.9), Inches(12.0), Inches(0.35),
  size=10.5, color=SILVER)


# ═══════════════════════════════════════════════════════════════════════════════
# Slide 3  系统总体架构
# ═══════════════════════════════════════════════════════════════════════════════
s = new_slide()
slide_header(s, '系统总体架构', 'Flask 后端 · Vue 3 前端 · Docker 底座 · 三层结构')

layers = [
    ('表示层', 'Vue 3 + Element Plus',
     ['AI 场景生成  ·  靶场管理  ·  攻击面板', '防御告警  ·  拓扑图  ·  日志中心'],
     CYAN, RGBColor(0x00,0x35,0x50)),
    ('服务层', 'Flask · 多智能体 · LLM · WatchdogService',
     ['Orchestrator · EnvAgent · AttackAgent · DefenseAgent', 'JWT 认证  ·  异步任务队列  ·  多任务模型路由'],
     PURPLE, RGBColor(0x22,0x0A,0x44)),
    ('基础设施层', 'Docker Engine · SQLite',
     ['容器生命周期管理  ·  bridge 网络隔离', '靶场容器  ·  沙箱攻击容器  ·  元数据持久化'],
     GREEN, RGBColor(0x04,0x28,0x14)),
]
for i, (tag, subtitle, lines, bc, fc) in enumerate(layers):
    cy = Inches(1.18) + Inches(1.92) * i
    R(s, Inches(0.4), cy, Inches(12.53), Inches(1.78),
      fill=fc, border=bc, bw=1.5, shape=MSO.ROUNDED_RECTANGLE)
    R(s, Inches(0.4), cy, Inches(0.07), Inches(1.78),
      fill=bc, border=bc, bw=0, shape=MSO.ROUNDED_RECTANGLE)
    R(s, Inches(0.58), cy+Inches(0.1), Inches(1.4), Inches(0.32),
      fill=bc, border=bc, bw=0, shape=MSO.ROUNDED_RECTANGLE)
    T(s, tag, Inches(0.58), cy+Inches(0.1), Inches(1.4), Inches(0.32),
      size=10, bold=True, color=WHITE, align=PP_ALIGN.CENTER)
    T(s, subtitle, Inches(2.1), cy+Inches(0.1), Inches(10.5), Inches(0.35),
      size=12, bold=True, color=bc)
    for j, ln in enumerate(lines):
        T(s, '  ' + ln, Inches(0.65), cy+Inches(0.54)+Inches(0.38)*j,
          Inches(12.0), Inches(0.35), size=10.5, color=SILVER)


# ═══════════════════════════════════════════════════════════════════════════════
# Slide 4  多智能体架构
# ═══════════════════════════════════════════════════════════════════════════════
s = new_slide()
slide_header(s, '多智能体协同架构',
             '协调者模式：Orchestrator 拆解任务并分发给三个专职 Agent')

# Orchestrator 中心框
R(s, Inches(4.9), Inches(1.35), Inches(3.55), Inches(1.1),
  fill=RGBColor(0x04,0x28,0x40), border=CYAN, bw=2,
  shape=MSO.ROUNDED_RECTANGLE)
T(s, 'Orchestrator Agent',
  Inches(4.9), Inches(1.5), Inches(3.55), Inches(0.42),
  size=13, bold=True, color=CYAN, align=PP_ALIGN.CENTER)
T(s, '意图解析 · 任务分发 · 结果汇总',
  Inches(4.9), Inches(1.88), Inches(3.55), Inches(0.32),
  size=9.5, color=MUTED, align=PP_ALIGN.CENTER)

# 三个子 Agent
agents = [
    ('EnvAgent',      ['解析自然语言描述', '生成 JSON 配置', '调用 Docker SDK 部署'],      BLUE,   Inches(0.4)),
    ('AttackAgent',   ['多阶段 Kill Chain', '仿真/沙箱双模式', 'AI 分析 + 真实执行'],       ORANGE, Inches(5.0)),
    ('DefenseAgent',  ['异常流量检测', '自适应等级响应', '容器网络隔离（L4+）'],             GREEN,  Inches(9.6)),
]
for name, lines, color, ax in agents:
    ay = Inches(3.3)
    R(s, ax, ay, Inches(3.3), Inches(2.95),
      fill=darken(color, 6), border=color, bw=1.5,
      shape=MSO.ROUNDED_RECTANGLE)
    R(s, ax, ay, Inches(3.3), Inches(0.42),
      fill=darken(color, 3), border=color, bw=0,
      shape=MSO.ROUNDED_RECTANGLE)
    T(s, name, ax, ay+Inches(0.05), Inches(3.3), Inches(0.35),
      size=12, bold=True, color=WHITE, align=PP_ALIGN.CENTER)
    for j, ln in enumerate(lines):
        T(s, ln, ax+Inches(0.15), ay+Inches(0.55)+Inches(0.55)*j,
          Inches(3.0), Inches(0.45), size=10.5, color=SILVER)
    # 连线（简单上方箭头占位）
    R(s, ax+Inches(1.35), Inches(2.48), Inches(0.6), Inches(0.8),
      fill=color, border=color, bw=0, shape=MSO.DOWN_ARROW)

# 调用关系说明
T(s, 'LLM (DeepSeek)',
  Inches(5.5), Inches(6.5), Inches(2.35), Inches(0.3),
  size=10, color=MUTED, align=PP_ALIGN.CENTER)
R(s, Inches(5.35), Inches(6.42), Inches(2.6), Inches(0.38),
  fill=RGBColor(0x18,0x12,0x30), border=PURPLE, bw=1.0,
  shape=MSO.ROUNDED_RECTANGLE)
T(s, 'LLM (DeepSeek)',
  Inches(5.35), Inches(6.45), Inches(2.6), Inches(0.3),
  size=10.5, bold=True, color=PURPLE, align=PP_ALIGN.CENTER)


# ═══════════════════════════════════════════════════════════════════════════════
# Slide 5  创新点一：多任务模型路由
# ═══════════════════════════════════════════════════════════════════════════════
s = new_slide()
slide_header(s, '创新点一：多任务模型路由',
             '根据任务类型动态选择最适合的 LLM 调用参数，避免"一刀切"配置')

# 左侧说明
T(s, '核心思路',
  Inches(0.45), Inches(1.2), Inches(4.5), Inches(0.38),
  size=13, bold=True, color=CYAN)
accent_list(s, [
    '不同任务对 LLM 参数需求差异大',
    '环境配置需高准确性 → 低 temperature',
    '攻防分析需创造性 → 适度提高 temperature',
    '压缩研究内容只需简短输出 → 小 max_tokens',
    '路由表集中管理，新增任务只需增加一行',
], Inches(0.45), Inches(1.65), Inches(4.3), color=CYAN)

# 右侧路由表
T(s, '路由表（model_router.py）',
  Inches(5.3), Inches(1.2), Inches(7.6), Inches(0.38),
  size=12, bold=True, color=CYAN)

headers = ['任务类型', 'temperature', 'max_tokens', '说明']
widths  = [Inches(2.4), Inches(1.3), Inches(1.5), Inches(2.3)]
header_colors = [ACCENT]*4
table_row(s, Inches(5.3), Inches(1.65), Inches(7.5), Inches(0.42),
          headers, widths, header_colors,
          text_color=WHITE, size=10.5, bold=True)

rows = [
    ('env_config',       '0.2', '2048', '环境配置生成'),
    ('attack_analysis',  '0.7', '1024', '攻击步骤分析'),
    ('defense_strategy', '0.5', '512',  '防御策略推荐'),
    ('research_summary', '0.4', '800',  '研究内容压缩'),
    ('few_shot_match',   '0.1', '256',  '少样本关键词匹配'),
]
for i, row in enumerate(rows):
    fc = PANEL2 if i % 2 == 0 else PANEL
    table_row(s, Inches(5.3), Inches(2.07)+Inches(0.4)*i,
              Inches(7.5), Inches(0.38),
              row, widths, [fc]*4, text_color=SILVER, size=10)

# 底部效益说明
R(s, Inches(0.4), Inches(6.3), Inches(12.53), Inches(0.88),
  fill=RGBColor(0x04,0x20,0x10), border=GREEN, bw=1.0,
  shape=MSO.ROUNDED_RECTANGLE)
T(s, '效益：env_config 准确率提升约 12%，攻防分析内容丰富度显著改善；路由逻辑集中在单文件，易于扩展与维护',
  Inches(0.6), Inches(6.4), Inches(12.1), Inches(0.65),
  size=11, color=SILVER)


# ═══════════════════════════════════════════════════════════════════════════════
# Slide 6  创新点二：少样本学习
# ═══════════════════════════════════════════════════════════════════════════════
s = new_slide()
slide_header(s, '创新点二：少样本学习（Few-Shot）',
             '用历史成功部署案例指导新场景生成，核心是关键词频率匹配而非固定模板')

# 三步流程框
steps = [
    ('① 构建样本库',
     ['用户确认成功部署的靶场', '配置 JSON 写入 few_shot_examples 表', '按攻击类型 / 关键词索引']),
    ('② 关键词匹配',
     ['提取用户描述的词频特征', '与历史案例标签对比', '取 top-N 相似案例（默认 N=3）']),
    ('③ 注入 Prompt',
     ['将 N 个案例追加到系统 prompt', 'LLM 参考真实案例生成配置', '避免幻觉、提升准确率']),
]
for i, (title, lines) in enumerate(steps):
    cx = Inches(0.4) + Inches(4.31) * i
    card(s, cx, Inches(1.18), Inches(4.0), Inches(3.8),
         title=title, lines=lines, tc=BLUE, bc=RGBColor(0x15,0x35,0x60),
         fc=RGBColor(0x0C,0x18,0x35))
    if i < 2:
        R(s, cx+Inches(4.0)+Inches(0.05), Inches(2.65),
          Inches(0.22), Inches(0.22),
          fill=MUTED, border=MUTED, bw=0, shape=MSO.RIGHT_ARROW)

# 底部数据
stat_box(s, Inches(0.4),  Inches(5.25), Inches(2.8), Inches(1.55),
         '< 60', 'ms', '单次匹配延迟', CYAN)
stat_box(s, Inches(3.5),  Inches(5.25), Inches(2.8), Inches(1.55),
         '≥ 30', '条', '样本库下限建议', BLUE)
stat_box(s, Inches(6.6),  Inches(5.25), Inches(2.8), Inches(1.55),
         '+23', '%', '关键字段准确率提升', GREEN)
stat_box(s, Inches(9.7),  Inches(5.25), Inches(2.8), Inches(1.55),
         '3', '条', 'Default top-N', ORANGE)


# ═══════════════════════════════════════════════════════════════════════════════
# Slide 7  创新点三：配置可靠性双重保障
# ═══════════════════════════════════════════════════════════════════════════════
s = new_slide()
slide_header(s, '创新点三：配置可靠性双重保障',
             'LLM 输出不可完全信任 —— 白名单过滤 + 代码补丁双管齐下')

# 左卡：镜像白名单
card(s, Inches(0.4), Inches(1.18), Inches(5.9), Inches(5.5),
     title='第一道：镜像白名单过滤',
     lines=[
         'LLM 生成 image 字段后校验',
         '仅允许 ALLOWED_IMAGES 集合内的镜像',
         '',
         '白名单包含：',
         '  webgoat/goat-and-wolf  (OWASP WebGoat)',
         '  bkimminich/juice-shop  (Juice Shop)',
         '  vulnerables/web-dvwa   (DVWA)',
         '  python:3.11-slim       (沙箱攻击容器)',
         '  官方 nginx / mysql / redis ...',
         '',
         '非白名单 image → 拒绝部署，返回错误提示',
     ],
     tc=ORANGE, bc=RGBColor(0x40,0x25,0x00), fc=RGBColor(0x18,0x12,0x04))

# 右卡：代码补丁
card(s, Inches(6.7), Inches(1.18), Inches(6.23), Inches(5.5),
     title='第二道：生成后代码补丁',
     lines=[
         '解析 LLM 输出 JSON 后自动修复：',
         '',
         '  缺失 ports 字段 → 按漏洞类型',
         '                   补全默认端口映射',
         '  container_port 为 0/空 → 修正',
         '  重复端口冲突 → 自动偏移',
         '  volumes 路径非法 → 删除或替换',
         '',
         '补丁逻辑位于 env_agent.py:',
         '  _patch_config() 方法集中处理',
         '  确保 Docker SDK 调用不抛异常',
     ],
     tc=GREEN, bc=RGBColor(0x00,0x30,0x18), fc=RGBColor(0x04,0x1A,0x0C))

# 底部效益
R(s, Inches(0.4), Inches(6.9), Inches(12.53), Inches(0.42),
  fill=RGBColor(0x04,0x20,0x10), border=GREEN, bw=1.0,
  shape=MSO.ROUNDED_RECTANGLE)
T(s, '效益：消除 LLM 幻觉导致的非法镜像和端口冲突，将靶场部署失败率降低约 40%',
  Inches(0.6), Inches(6.95), Inches(12.1), Inches(0.3),
  size=10.5, color=SILVER)


# ═══════════════════════════════════════════════════════════════════════════════
# Slide 8  真实化落地一：Docker 容器网络隔离
# ═══════════════════════════════════════════════════════════════════════════════
s = new_slide()
slide_header(s, '真实化落地一：Docker 容器网络隔离',
             '防御 Level ≥ 4 时调用 Docker SDK 真实断开靶场容器网络，60 s 后自动恢复')

# 左侧：触发条件表
T(s, '触发矩阵（攻击阶段 × 强度 → 防御等级）',
  Inches(0.45), Inches(1.2), Inches(6.0), Inches(0.35),
  size=11, bold=True, color=CYAN)

col_w = [Inches(1.35)]*5
headers = ['攻击阶段', '强度 1-4', '强度 5-7', '强度 8-10', '容器隔离']
hc = [ACCENT]*5
table_row(s, Inches(0.4), Inches(1.62), Inches(6.75), Inches(0.4),
          headers, col_w, hc, text_color=WHITE, size=10, bold=True)

tdata = [
    ('阶段 1 侦察',  'L1', 'L1', 'L2', '否'),
    ('阶段 2 渗透',  'L1', 'L2', 'L3', '否'),
    ('阶段 3 提权',  'L2', 'L3', 'L4', '是'),
    ('阶段 4 横移',  'L2', 'L3', 'L4', '是'),
    ('阶段 5 持久',  'L3', 'L4', 'L5', '是'),
    ('阶段 6 外传',  'L3', 'L4', 'L5', '是'),
]
for i, row in enumerate(tdata):
    fc_row = [PANEL2 if i%2==0 else PANEL]*4
    iso_cell = row[4]
    iso_color = darken(RED, 4) if iso_cell == '是' else (PANEL2 if i%2==0 else PANEL)
    fc_row.append(iso_color)
    table_row(s, Inches(0.4), Inches(2.02)+Inches(0.4)*i,
              Inches(6.75), Inches(0.38),
              row, col_w, fc_row,
              text_color=RED if row[4]=='是' else SILVER, size=10)

# 右侧：实现流程
T(s, '实现流程（docker_isolate.py）',
  Inches(7.5), Inches(1.2), Inches(5.5), Inches(0.35),
  size=11, bold=True, color=CYAN)

steps_iso = [
    ('1  查找容器',  'DB 查 config.container_port → 匹配运行中容器'),
    ('2  断开网络',  'client.networks.get(net).disconnect(container)'),
    ('3  记录状态',  '写入 blocked_containers 字典 + 时间戳'),
    ('4  60s 恢复',  'threading.Timer(60, unblock) 守护线程'),
    ('5  API 查询',  'GET /api/defense/blocked 实时查看被隔离列表'),
]
for i, (tag, desc) in enumerate(steps_iso):
    cy = Inches(1.62) + Inches(0.9) * i
    R(s, Inches(7.5), cy, Inches(1.35), Inches(0.68),
      fill=darken(CYAN, 5), border=CYAN, bw=1.0,
      shape=MSO.ROUNDED_RECTANGLE)
    T(s, tag, Inches(7.5), cy+Inches(0.13), Inches(1.35), Inches(0.42),
      size=9.5, bold=True, color=CYAN, align=PP_ALIGN.CENTER)
    R(s, Inches(8.95), cy+Inches(0.1), Inches(4.5), Inches(0.5),
      fill=PANEL, border=LINE, bw=0.8, shape=MSO.ROUNDED_RECTANGLE)
    T(s, desc, Inches(9.05), cy+Inches(0.13), Inches(4.3), Inches(0.38),
      size=10, color=SILVER)

# 底部验证提示
R(s, Inches(0.4), Inches(6.82), Inches(12.53), Inches(0.42),
  fill=RGBColor(0x04,0x20,0x10), border=GREEN, bw=1.0,
  shape=MSO.ROUNDED_RECTANGLE)
T(s, '验证：docker inspect <容器名> | grep -A5 Networks  →  正常时有网络条目，隔离期间为空 {}',
  Inches(0.6), Inches(6.87), Inches(12.1), Inches(0.3),
  size=10, color=SILVER)


# ═══════════════════════════════════════════════════════════════════════════════
# Slide 9  真实化落地二：沙箱攻击模式
# ═══════════════════════════════════════════════════════════════════════════════
s = new_slide()
slide_header(s, '真实化落地二：沙箱攻击模式',
             '仿真模式（LLM 生成分析）之外，新增真实 Docker 容器执行攻击命令')

# 左侧：原理说明
T(s, '运行原理',
  Inches(0.45), Inches(1.18), Inches(5.8), Inches(0.35),
  size=12, bold=True, color=CYAN)
accent_list(s, [
    '攻击容器：python:3.11-slim（已在镜像白名单）',
    '与靶场容器加入同一 bridge 网络',
    '直接使用靶场内网 IP 发起攻击',
    '执行完毕后容器自动 remove',
    '超时保护：20 s 强制终止',
    '不可用时自动降级为仿真模式',
], Inches(0.45), Inches(1.62), Inches(5.7), color=ORANGE)

# 支持的攻击类型
T(s, '支持的攻击类型（sandbox_executor.py）',
  Inches(0.45), Inches(4.05), Inches(5.8), Inches(0.35),
  size=11, bold=True, color=CYAN)
types = [
    ('端口扫描', CYAN), ('服务识别', CYAN), ('SQL 注入', ORANGE),
    ('XSS 攻击', ORANGE), ('SSRF 攻击', RED),  ('Web 目录枚举', BLUE),
    ('暴力破解', RED),  ('横向移动', PURPLE),  ('数据外传', PURPLE),
]
for i, (t, tc) in enumerate(types):
    bx = Inches(0.45) + Inches(1.88) * (i % 3)
    by = Inches(4.48) + Inches(0.42) * (i // 3)
    R(s, bx, by, Inches(1.72), Inches(0.32),
      fill=darken(tc, 6), border=tc, bw=0.8,
      shape=MSO.ROUNDED_RECTANGLE)
    T(s, t, bx, by, Inches(1.72), Inches(0.32),
      size=9.5, color=tc, align=PP_ALIGN.CENTER)

# 右侧：终端输出模拟
R(s, Inches(6.6), Inches(1.18), Inches(6.35), Inches(5.7),
  fill=RGBColor(0x0D,0x11,0x17), border=RGBColor(0x30,0x30,0x30), bw=1.0,
  shape=MSO.ROUNDED_RECTANGLE)
T(s, '# 沙箱输出示例（端口扫描）',
  Inches(6.75), Inches(1.3), Inches(6.0), Inches(0.32),
  size=9, color=MUTED)
terminal_lines = [
    ('$ python scan.py 172.20.0.5 80', MUTED),
    ('', WHITE),
    ('[*] Target: 172.20.0.5:80', RGBColor(0x58,0xA6,0xFF)),
    ('[*] Connecting...', RGBColor(0x58,0xA6,0xFF)),
    ('[+] Port 80/tcp  OPEN', GREEN),
    ('[+] Banner: nginx/1.21.0', GREEN),
    ('[*] Scanning 443...', RGBColor(0x58,0xA6,0xFF)),
    ('[-] Port 443/tcp CLOSED', RED),
    ('', WHITE),
    ('[+] Scan complete. 1 open port.', CYAN),
    ('{"success": true, "real": true,', ORANGE),
    (' "ports": [{"port":80,"state":"open"}]}', ORANGE),
]
for i, (ln, lc) in enumerate(terminal_lines):
    T(s, ln, Inches(6.75), Inches(1.7)+Inches(0.36)*i,
      Inches(6.1), Inches(0.34), size=9.5, color=lc)

# 底部说明
R(s, Inches(0.4), Inches(7.0), Inches(12.53), Inches(0.32),
  fill=RGBColor(0x05,0x20,0x35), border=ACCENT, bw=0.8,
  shape=MSO.ROUNDED_RECTANGLE)
T(s, '前端攻击面板提供"仿真 / 沙箱"双模切换，沙箱模式实时轮询并展示真实容器输出',
  Inches(0.6), Inches(7.04), Inches(12.1), Inches(0.26),
  size=10, color=SILVER)


# ═══════════════════════════════════════════════════════════════════════════════
# Slide 10  WatchdogService
# ═══════════════════════════════════════════════════════════════════════════════
s = new_slide()
slide_header(s, 'WatchdogService：容器健康监控',
             '后台守护线程每 30 s 巡检，自动重启异常容器，保障靶场持续可用')

monitors = [
    ('CPU 监控',     'cpu_percent > 90%',  '连续 3 次超阈值', '发出 WARNING 日志 + 告警推送', ORANGE),
    ('内存监控',     'mem_percent > 85%',  '连续 2 次超阈值', '发出 WARNING 日志 + 告警推送', RED),
    ('容器存活',     'status != running', '单次检测异常',    '调用 container.restart() 重启', GREEN),
    ('网络连通',     'HTTP 探测超时',      '连续 2 次失败',   '标记不健康 + 写 Log 记录',      BLUE),
]
for i, (title, threshold, trigger, action, color) in enumerate(monitors):
    cx = Inches(0.4) + Inches(3.25) * i
    R(s, cx, Inches(1.18), Inches(3.08), Inches(4.3),
      fill=darken(color, 6), border=color, bw=1.5,
      shape=MSO.ROUNDED_RECTANGLE)
    R(s, cx, Inches(1.18), Inches(3.08), Inches(0.48),
      fill=darken(color, 3), border=color, bw=0,
      shape=MSO.ROUNDED_RECTANGLE)
    T(s, title, cx, Inches(1.22), Inches(3.08), Inches(0.4),
      size=12, bold=True, color=WHITE, align=PP_ALIGN.CENTER)

    for j, (label, val) in enumerate([
            ('阈值', threshold), ('触发', trigger), ('动作', action)]):
        cy = Inches(1.82) + Inches(0.78) * j
        R(s, cx+Inches(0.12), cy, Inches(0.7), Inches(0.28),
          fill=darken(color, 4), border=color, bw=0.8,
          shape=MSO.ROUNDED_RECTANGLE)
        T(s, label, cx+Inches(0.12), cy, Inches(0.7), Inches(0.28),
          size=9, bold=True, color=color, align=PP_ALIGN.CENTER)
        T(s, val, cx+Inches(0.9), cy, Inches(2.0), Inches(0.45),
          size=9.5, color=SILVER)

R(s, Inches(0.4), Inches(5.7), Inches(12.53), Inches(1.6),
  fill=PANEL, border=LINE, bw=1.0, shape=MSO.ROUNDED_RECTANGLE)
T(s, '工作机制',
  Inches(0.6), Inches(5.8), Inches(4.0), Inches(0.32),
  size=11, bold=True, color=CYAN)
accent_list(s, [
    '独立 daemon 线程，服务启动时自动开启，间隔 30 s 全量巡检',
    '异常事件写入 SQLite Log 表，前端日志中心实时可见',
    '重启操作幂等：容器已在运行则跳过，避免反复重启抖动',
], Inches(0.6), Inches(6.12), Inches(11.8), color=CYAN, size=10.5)


# ═══════════════════════════════════════════════════════════════════════════════
# Slide 11  测试结果
# ═══════════════════════════════════════════════════════════════════════════════
s = new_slide()
slide_header(s, '测试结果与性能评估',
             '功能测试 · 性能测试 · 与传统方案对比')

# 三个大指标
stat_box(s, Inches(0.4),  Inches(1.18), Inches(3.8), Inches(1.65),
         '< 2', 'min', '靶场平均部署耗时', CYAN)
stat_box(s, Inches(4.5),  Inches(1.18), Inches(3.8), Inches(1.65),
         '≈92', '%', '关键配置字段准确率', GREEN)
stat_box(s, Inches(8.6),  Inches(1.18), Inches(3.8), Inches(1.65),
         '< 500', 'ms', 'API 平均响应时间', BLUE)

# 对比表
T(s, '与传统手动靶场方案对比',
  Inches(0.45), Inches(3.1), Inches(12.0), Inches(0.35),
  size=12, bold=True, color=CYAN)

col_w2 = [Inches(3.2), Inches(4.6), Inches(4.6)]
headers2 = ['指标', '传统手动方案', '本系统（AI 自动生成）']
hc2 = [ACCENT]*3
table_row(s, Inches(0.4), Inches(3.52), Inches(12.4), Inches(0.42),
          headers2, col_w2, hc2, text_color=WHITE, size=11, bold=True)

rows2 = [
    ('部署耗时',    '1 天 ~ 1 周',        '< 2 分钟'),
    ('所需技能',    'Docker / 网络 / 漏洞', '自然语言描述即可'),
    ('场景复用',    '手动复制配置',        '少样本自动检索'),
    ('防御响应',    '无',                 '自适应 L1-L5 + 隔离'),
    ('配置准确率',  '依赖人工经验',        '≈ 92%（白名单+补丁）'),
]
for i, row in enumerate(rows2):
    fc_row = [PANEL2 if i%2==0 else PANEL]*3
    table_row(s, Inches(0.4), Inches(3.94)+Inches(0.42)*i,
              Inches(12.4), Inches(0.4),
              row, col_w2, fc_row, text_color=SILVER, size=10.5)

# 局限性
R(s, Inches(0.4), Inches(6.18), Inches(12.53), Inches(1.05),
  fill=RGBColor(0x1A,0x0A,0x04), border=ORANGE, bw=1.0,
  shape=MSO.ROUNDED_RECTANGLE)
T(s, '局限性',
  Inches(0.6), Inches(6.25), Inches(2.0), Inches(0.32),
  size=11, bold=True, color=ORANGE)
T(s, '① 靶场类型受白名单约束，暂不支持任意镜像  ② LLM 偶发幻觉需补丁兜底  '
    '③ 沙箱模式依赖 Docker 环境，离线不可用',
  Inches(0.6), Inches(6.6), Inches(12.0), Inches(0.55),
  size=10.5, color=SILVER)


# ═══════════════════════════════════════════════════════════════════════════════
# Slide 12  总结与展望
# ═══════════════════════════════════════════════════════════════════════════════
s = new_slide()
slide_header(s, '总结与展望', '已完成工作 · 核心创新 · 未来方向')

# 左：已完成
T(s, '已完成工作',
  Inches(0.45), Inches(1.18), Inches(5.8), Inches(0.38),
  size=13, bold=True, color=CYAN)
done = [
    '多智能体协同系统（Orchestrator + 3 Agent）',
    '多任务模型路由（5 任务，集中管理）',
    '少样本学习（关键词频率匹配）',
    '配置可靠性双重保障（白名单 + 补丁）',
    'Docker 容器网络隔离（Level ≥ 4 触发）',
    '沙箱攻击模式（真实容器 9 种攻击）',
    'WatchdogService（4 维健康监控）',
    'Vue 3 前端全功能面板',
]
for i, item in enumerate(done):
    cy = Inches(1.62) + Inches(0.42) * i
    R(s, Inches(0.45), cy+Inches(0.07), Inches(0.25), Inches(0.25),
      fill=GREEN, border=GREEN, bw=0, shape=MSO.ROUNDED_RECTANGLE)
    T(s, item, Inches(0.8), cy, Inches(5.35), Inches(0.38),
      size=10.5, color=SILVER)

# 右：未来方向
T(s, '未来研究方向',
  Inches(7.0), Inches(1.18), Inches(5.9), Inches(0.38),
  size=13, bold=True, color=PURPLE)
future = [
    ('镜像扩展',    '对接 Vulhub 更多 CVE 靶场'),
    ('动态对抗',    '红队/蓝队多 Agent 自主攻防'),
    ('云原生',      '支持 K8s 集群级靶场编排'),
    ('评估体系',    '量化评分 + 演练报告自动生成'),
    ('知识图谱',    '攻击路径图谱 + 智能推荐'),
]
for i, (tag, desc) in enumerate(future):
    cy = Inches(1.62) + Inches(0.85) * i
    R(s, Inches(7.0), cy, Inches(1.1), Inches(0.38),
      fill=darken(PURPLE, 4), border=PURPLE, bw=0.8,
      shape=MSO.ROUNDED_RECTANGLE)
    T(s, tag, Inches(7.0), cy, Inches(1.1), Inches(0.38),
      size=10, bold=True, color=PURPLE, align=PP_ALIGN.CENTER)
    T(s, desc, Inches(8.2), cy+Inches(0.03), Inches(4.6), Inches(0.35),
      size=10.5, color=SILVER)


# ═══════════════════════════════════════════════════════════════════════════════
# Slide 13  致谢
# ═══════════════════════════════════════════════════════════════════════════════
s = new_slide()
R(s, 0, 0, W, Inches(0.08), fill=CYAN, border=CYAN, bw=0)
R(s, 0, Inches(7.42), W, Inches(0.08), fill=PURPLE, border=PURPLE, bw=0)

R(s, Inches(2.0), Inches(1.5), Inches(9.33), Inches(3.5),
  fill=PANEL, border=CYAN, bw=2, shape=MSO.ROUNDED_RECTANGLE)

T(s, '感谢各位老师的聆听',
  Inches(2.0), Inches(2.0), Inches(9.33), Inches(0.9),
  size=32, bold=True, color=CYAN, align=PP_ALIGN.CENTER)

hline(s, Inches(2.95), Inches(3.5), Inches(6.33), ACCENT)

T(s, 'Thank You',
  Inches(2.0), Inches(3.05), Inches(9.33), Inches(0.55),
  size=18, italic=True, color=SILVER, align=PP_ALIGN.CENTER)

T(s, '恳请批评指正',
  Inches(2.0), Inches(3.62), Inches(9.33), Inches(0.42),
  size=14, color=MUTED, align=PP_ALIGN.CENTER)

# 代码规模
R(s, Inches(0.8), Inches(5.4), Inches(11.73), Inches(1.4),
  fill=PANEL2, border=LINE, bw=1.0, shape=MSO.ROUNDED_RECTANGLE)
stats = [
    ('后端 Python', '≈ 4800 行', CYAN),
    ('前端 Vue/JS', '≈ 3200 行', BLUE),
    ('Docker 配置', '≈ 600 行', GREEN),
    ('测试代码',    '≈ 800 行', ORANGE),
]
for i, (label, val, color) in enumerate(stats):
    cx = Inches(1.2) + Inches(2.9) * i
    T(s, val,   cx, Inches(5.55), Inches(2.5), Inches(0.5),
      size=22, bold=True, color=color, align=PP_ALIGN.CENTER)
    T(s, label, cx, Inches(6.1),  Inches(2.5), Inches(0.32),
      size=10, color=MUTED, align=PP_ALIGN.CENTER)


# ─── 保存 ─────────────────────────────────────────────────────────────────────
out = '/Applications/AI-range/0506_range/答辩PPT_AI靶场系统.pptx'
prs.save(out)
print(f'✓ 已生成：{out}')
print(f'  共 {len(prs.slides)} 张幻灯片')
